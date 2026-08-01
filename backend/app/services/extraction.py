from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

def extract_text(file_type: str, file_path: str):
    if file_type == "pdf":
        try:
            pdf_reader = PdfReader(file_path)
        except Exception as e:
            raise ValueError(f"Could not open PDF file: {e}")

        if pdf_reader.is_encrypted:
            raise ValueError("PDF is password-protected and cannot be processed!")
        
        extracted_text = []
        nr_pages = len(pdf_reader.pages)

        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                extracted_text.append(text)

        full_text = "\n".join(extracted_text)

    elif file_type == "docx":
        try:
            doc = DocxDocument(file_path)
        except Exception as e:
            raise ValueError(f"Could not open DOCX file: {e}")
        
        extracted_text = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                extracted_text.append(text)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        extracted_text.append(text)

        full_text = "\n".join(extracted_text)

        # Estimate the number of pages
        if len(full_text) > 0:
            nr_pages = max(1, len(full_text) // 1500)
        else:
            nr_pages = 0

    elif file_type == "pptx":
        try:
            pres = PptxPresentation(file_path)
        except Exception as e:
            raise ValueError(f"Could not open PPTX file: {e}")
        
        extracted_text = []
        nr_pages = len(pres.slides)

        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        extracted_text.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                extracted_text.append(text)

        full_text = "\n".join(extracted_text)

    else:
        raise ValueError("Unsupported file type!")

    if not full_text.strip():
        raise ValueError("Extraction returned empty text or failed!")

    return full_text, nr_pages